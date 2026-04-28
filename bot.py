import os
import logging
import requests
from aiogram import Bot, Dispatcher, types
from aiogram.types import ReplyKeyboardMarkup, KeyboardButton
from aiogram.utils import executor
from pptx import Presentation
from pptx.util import Inches
from openai import OpenAI
import random

# ====== CONFIG ======
BOT_TOKEN = os.getenv("BOT_TOKEN")
import os

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
CHANNEL_USERNAME = "@master_botuz"
PRICE_PER_SLIDE = 400

bot = Bot(token=BOT_TOKEN)
dp = Dispatcher(bot)
client = OpenAI(api_key=OPENAI_API_KEY)

logging.basicConfig(level=logging.INFO)

users = {}

# ====== SUBSCRIBE CHECK ======
async def check_sub(user_id):
    try:
        member = await bot.get_chat_member(CHANNEL_USERNAME, user_id)
        return member.status in ["member", "creator", "administrator"]
    except:
        return False

# ====== START ======
@dp.message_handler(commands=['start'])
async def start(msg: types.Message):
    if not await check_sub(msg.from_user.id):
        kb = ReplyKeyboardMarkup(resize_keyboard=True)
        kb.add("✅ Tekshirish")
        await msg.answer(
            f"Kanalga obuna bo‘ling:\nhttps://t.me/master_botuz\n\nSo‘ng tekshiring",
            reply_markup=kb
        )
        return

    await msg.answer("Mustaqil ish mavzusini yuboring:")
    users[msg.from_user.id] = {}

# ====== CHECK BUTTON ======
@dp.message_handler(lambda m: m.text == "✅ Tekshirish")
async def check(msg: types.Message):
    if await check_sub(msg.from_user.id):
        await msg.answer("Mavzuni yuboring:")
    else:
        await msg.answer("Hali obuna bo‘lmadingiz!")

# ====== TOPIC ======
@dp.message_handler(lambda m: m.from_user.id in users and "topic" not in users[m.from_user.id])
async def get_topic(msg: types.Message):
    users[msg.from_user.id]["topic"] = msg.text

    kb = ReplyKeyboardMarkup(resize_keyboard=True)
    kb.add("Minimal", "Dark", "Creative")

    await msg.answer("Dizayn tanlang:", reply_markup=kb)

# ====== DESIGN ======
@dp.message_handler(lambda m: m.text in ["Minimal", "Dark", "Creative"])
async def design(msg: types.Message):
    users[msg.from_user.id]["design"] = msg.text
    await msg.answer("Nechta slayd kerak?")

# ====== SLIDE COUNT ======
@dp.message_handler(lambda m: m.text.isdigit())
async def slide_count(msg: types.Message):
    count = int(msg.text)
    users[msg.from_user.id]["count"] = count

    price = count * PRICE_PER_SLIDE
    users[msg.from_user.id]["price"] = price

    await msg.answer(f"Narx: {price} so‘m\n\nTo‘lov qilgandan keyin 'to‘ladim' deb yozing")

# ====== PAYMENT FAKE ======
@dp.message_handler(lambda m: "to‘ladim" in m.text.lower())
async def paid(msg: types.Message):
    data = users[msg.from_user.id]

    await msg.answer("Yuklanmoqda...")

    file = create_ppt(
        data["topic"],
        data["design"],
        data["count"]
    )

    await msg.answer_document(open(file, "rb"))

# ====== AI TEXT ======
def generate_text(topic):
    prompt = f"""
    {topic} mavzusida slayd uchun professional, qisqa va tushunarli matn yoz.
    Oddiy emas, kitob uslubida emas, prezentatsiya uslubida yoz.
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4.1-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        return response.choices[0].message.content.strip()
    except:
        return f"{topic} haqida muhim ma'lumot"

# ====== IMAGE ======
def get_image(topic):
    url = f"https://source.unsplash.com/800x600/?{topic}"
    img_data = requests.get(url).content
    file = f"{topic}.jpg"
    with open(file, "wb") as f:
        f.write(img_data)
    return file

# ====== PPT ======
def create_ppt(topic, design, count):
    prs = Presentation()

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = topic
    slide.placeholders[1].text = "Mustaqil ish"

    for i in range(count):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{topic}"
        slide.placeholders[1].text = generate_text(topic)

    # images
    for i in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        img = get_image(topic)
        slide.shapes.add_picture(img, Inches(1), Inches(1.5), width=Inches(6))

    # thanks
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "E’tiboringiz uchun rahmat!"

    file = f"{topic}.pptx"
    prs.save(file)

    return file

# ====== RUN ======
if __name__ == "__main__":
    executor.start_polling(dp)
